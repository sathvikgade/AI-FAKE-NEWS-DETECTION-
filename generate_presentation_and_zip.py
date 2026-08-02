import os
import zipfile
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # 16:9 Aspect Ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]
    
    # Color Palette
    DARK_BLUE = RGBColor(15, 23, 42)
    NAVY_BLUE = RGBColor(30, 58, 138)
    ACCENT_BLUE = RGBColor(59, 130, 246)
    TEXT_DARK = RGBColor(30, 41, 59)
    TEXT_MUTED = RGBColor(100, 116, 139)
    WHITE = RGBColor(255, 255, 255)
    GREEN = RGBColor(16, 185, 129)

    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="AI & ML Summer Internship Program 2026"):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.name = 'Arial'
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.name = 'Arial'
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = NAVY_BLUE

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide1, DARK_BLUE)

    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "PROJECT 1: SUMMER INTERNSHIP PROGRAM IN AI & ML 2026"
    p0.font.name = 'Arial'
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_BLUE
    p0.space_after = Pt(20)

    p1 = tf.add_paragraph()
    p1.text = "AI-Powered Fake News Detection\nUsing Text Classification"
    p1.font.name = 'Arial'
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(25)

    p2 = tf.add_paragraph()
    p2.text = "Indian Institute of Computing and Technology (IICT)\nAffiliated: I-STEM, Office of the Principal Scientific Adviser to the Government of India"
    p2.font.name = 'Arial'
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(148, 163, 184)

    # SLIDE 2: Problem Statement & Objectives
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide2, "Problem Statement & Internship Objectives")
    
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    bullets = [
        ("Problem Statement", "Build an end-to-end machine learning pipeline from scratch to classify news articles as Real or Fake without relying on pre-built turn-key wrappers."),
        ("Motivation", "Digital media proliferation enables fake news to spread exponentially. Automated NLP models offer scalable, real-time verification to prevent public deception."),
        ("30-Day Workflow", "Structured across 4 weeks: Data Collection & Cleaning (Week 1), Feature Engineering & TF-IDF (Week 2), Model Building (Week 3), Evaluation & IEEE Report (Week 4)."),
        ("Key Requirements", "Implement 4 distinct models: KNN (Non-Parametric), Logistic Regression (Parametric), Random Forest (Ensemble), and Simple Neural Network (MLP).")
    ]

    for title, desc in bullets:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {title}: "
        r1.font.bold = True
        r1.font.size = Pt(16)
        r1.font.color.rgb = NAVY_BLUE
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(14)

    # SLIDE 3: Dataset Description
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide3, "Dataset Overview & Class Distribution")

    txBox = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    ds_bullets = [
        ("Dataset Source", "Kaggle Fake and Real News Dataset (comprising 39,999 total news articles)."),
        ("Class Balance", "Fake News: 20,886 articles (52.2%) | Real News: 19,113 articles (47.8%)."),
        ("Attributes Extracted", "index, title (headline), text (body), subject, date, and label (target)."),
        ("Key Structural Findings", "Real articles frequently contain news agency markers (e.g. 'WASHINGTON (Reuters) -'), while Fake articles display all-caps clickbait titles and emotional rhetoric.")
    ]

    for title, desc in ds_bullets:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {title}: "
        r1.font.bold = True
        r1.font.size = Pt(16)
        r1.font.color.rgb = NAVY_BLUE
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(14)

    # SLIDE 4: Week 1 & 2 - Preprocessing & Feature Engineering
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide4, "Week 1 & 2: Text Cleaning & TF-IDF Vectorization")

    txBox = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    steps = [
        ("Agency Normalization", "Removed static agency prefixes (e.g., 'WASHINGTON (Reuters) -') to prevent shortcut bias."),
        ("Text Sanitization", "Converted text to lowercase and stripped punctuation/special characters using Regex pattern [^a-z\\s]."),
        ("Manual Tokenization", "Split text streams into clean word tokens and removed non-informative English stop-words via NLTK."),
        ("TF-IDF Vectorization", "Extracted Unigrams + Bigrams (10,000 max features) using Sublinear TF scaling: TF-IDF(t, d) = TF(t, d) * IDF(t, D).")
    ]

    for title, desc in steps:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {title}: "
        r1.font.bold = True
        r1.font.size = Pt(16)
        r1.font.color.rgb = NAVY_BLUE
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(14)

    # SLIDE 5: Week 3 - Machine Learning Models Implemented
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide5, "Week 3: Algorithm Implementation Paradigms")

    txBox = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    algos = [
        ("K-Nearest Neighbors (KNN)", "Non-parametric distance-based classifier evaluating k=5 nearest neighbors via Euclidean/Cosine distance."),
        ("Logistic Regression", "Parametric linear classifier modeling log-odds using sigmoid activation function."),
        ("Random Forest Classifier", "Ensemble bagging model combining 100 de-correlated decision trees with random feature subsampling."),
        ("Simple Neural Network (MLP)", "Multi-Layer Perceptron containing 100 hidden neurons, ReLU activation, and Adam optimizer.")
    ]

    for title, desc in algos:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {title}: "
        r1.font.bold = True
        r1.font.size = Pt(16)
        r1.font.color.rgb = NAVY_BLUE
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(14)

    # SLIDE 6: Week 4 - Quantitative Results & IEEE Matrix
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide6, "Week 4: Experimental Results & IEEE Matrix")

    rows, cols = 5, 5
    table_shape = slide6.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5))
    table = table_shape.table

    headers = ["Algorithm", "Model Type", "Accuracy", "Precision", "F1-Score"]
    for col_idx, h in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(15)
        p.runs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_BLUE

    matrix_data = [
        ["KNN", "Non-Parametric", "76.46%", "91.00%", "70.00%"],
        ["Logistic Regression", "Parametric Linear", "98.67%", "98.00%", "99.00%"],
        ["Neural Network (MLP)", "Deep Learning", "99.09%", "99.00%", "99.00%"],
        ["Random Forest (Winner)", "Ensemble Trees", "99.72%", "100.00%", "100.00%"]
    ]

    for row_idx, row_data in enumerate(matrix_data, start=1):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(14)
            p.runs[0].font.name = 'Arial'
            if row_idx == 4:
                p.runs[0].font.bold = True
                p.runs[0].font.color.rgb = GREEN

    # SLIDE 7: Discussion - Parametric vs. Non-Parametric
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide7, "Discussion: Model Performance Trade-Offs")

    txBox = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    disc = [
        ("Parametric Efficiency", "Logistic Regression (98.67%) trained in < 2 seconds, proving linear decision boundaries thrive in high-dimensional sparse TF-IDF text vector spaces."),
        ("Non-Parametric Degradation", "KNN achieved the lowest accuracy (76.46%) due to the 'Curse of Dimensionality', where Euclidean distance metrics become uniform and noisy in 5,000D space."),
        ("Ensemble Dominance", "Random Forest achieved top accuracy (99.72%) because combining 100 decision trees effectively cancels out individual tree variance and prevents term overfitting."),
        ("Verification Test Suite", "Achieved 100% verification accuracy (30/30 passed) on sampled Real vs Fake dataset test cases.")
    ]

    for title, desc in disc:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {title}: "
        r1.font.bold = True
        r1.font.size = Pt(16)
        r1.font.color.rgb = NAVY_BLUE
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(14)

    # SLIDE 8: Conclusion & Future Scope
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide8, "Conclusion & Future Work")

    txBox = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    conc = [
        ("Core Takeaway", "Text preprocessing combined with TF-IDF vectorization and Random Forest ensemble learning delivers state-of-the-art fake news classification."),
        ("Deployments Built", "Built a full modular Python codebase, a real-time prediction CLI, and an interactive web dashboard (http://localhost:3000)."),
        ("Project Limitations", "Potential performance decay on newly emerging topics outside the 2016-2017 political training corpus."),
        ("Future Scope", "Integration of Pre-trained Transformer Architectures (BERT, RoBERTa, LLMs) and multimodal graph propagation tracking.")
    ]

    for title, desc in conc:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {title}: "
        r1.font.bold = True
        r1.font.size = Pt(16)
        r1.font.color.rgb = NAVY_BLUE
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(14)

    output_pptx = r"C:\Users\LENOVO\OneDrive\Desktop\AI_FAKE_NEWS_PROJECT\AI_Fake_News_Detection_Presentation.pptx"
    prs.save(output_pptx)
    print(f"Presentation saved successfully to: {output_pptx}")
    return output_pptx

def create_project_zip():
    base_dir = r"C:\Users\LENOVO\OneDrive\Desktop\AI_FAKE_NEWS_PROJECT"
    zip_filename = os.path.join(base_dir, "AI_Fake_News_Detection_Project.zip")

    files_to_zip = [
        ("AI_Fake_News_Detection_Presentation.pptx", "Presentation/AI_Fake_News_Detection_Presentation.pptx"),
        ("IEEE_Fake_News_Detection_Report.docx", "Documentation/IEEE_Fake_News_Detection_Report.docx"),
        ("IEEE_Fake_News_Detection_Report.md", "Documentation/IEEE_Fake_News_Detection_Report.md"),
        ("main.py", "Source_Code/main.py"),
        ("preprocessing.py", "Source_Code/preprocessing.py"),
        ("feature_engineering.py", "Source_Code/feature_engineering.py"),
        ("train.py", "Source_Code/train.py"),
        ("evaluate.py", "Source_Code/evaluate.py"),
        ("predict.py", "Source_Code/predict.py"),
        ("fake_news_pipeline.py", "Source_Code/fake_news_pipeline.py"),
        ("fake_news_detection.ipynb", "Source_Code/fake_news_detection.ipynb"),
        ("test_dataset_predictions.py", "Source_Code/test_dataset_predictions.py"),
        ("index.html", "Web_Dashboard/index.html"),
        ("style.css", "Web_Dashboard/style.css"),
        ("app.js", "Web_Dashboard/app.js"),
        ("package.json", "Web_Dashboard/package.json"),
        ("README.md", "README.md")
    ]

    with zipfile.ZipFile(zip_filename, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for file_name, arcname in files_to_zip:
            file_path = os.path.join(base_dir, file_name)
            if os.path.exists(file_path):
                zipf.write(file_path, arcname)
                print(f"Added to ZIP: {arcname}")

    print(f"\nZIP archive created successfully at: {zip_filename}")
    return zip_filename

if __name__ == "__main__":
    create_presentation()
    create_project_zip()
